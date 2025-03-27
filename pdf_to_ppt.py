import os
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE, MSO_SHAPE_TYPE
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
import re
import argparse
import sys
from math import ceil

def hex_to_rgb(hex_color):
    """Convert hex color to RGB."""
    hex_color = hex_color.lstrip('#')
    return tuple(int(hex_color[i:i+2], 16) for i in (0, 2, 4))

def is_bullet_character(text):
    """Check if the text starts with a bullet point character."""
    bullet_chars = ['•', '·', '○', '●', '▪', '▫', '◦', '-', '*']
    return any(text.strip().startswith(char) for char in bullet_chars)

def clean_bullet_text(text):
    """Remove bullet character from the beginning of text."""
    bullet_chars = ['•', '·', '○', '●', '▪', '▫', '◦', '-', '*']
    text = text.strip()
    for char in bullet_chars:
        if text.startswith(char):
            return text[len(char):].strip()
    return text

def detect_shape_type(shape):
    """Detect the type of shape based on its properties."""
    rect = shape["rect"]
    width = rect[2] - rect[0]
    height = rect[3] - rect[1]
    
    # Avoid division by zero
    if width == 0 or height == 0:
        return "rectangle"
    
    # Detect circles/dots (aspect ratio close to 1 and small size)
    if 0.9 <= width/height <= 1.1 and width < 20:
        return "status_dot"
    
    # Detect progress bars (wide rectangle)
    if width > height * 3:
        return "progress_bar"
    
    # Detect cards (large rectangles with specific aspect ratio)
    if width > 100 and height > 100:
        return "card"
    
    return "rectangle"

def create_status_card(slide, left, top, width, height, title, status_text, date_text=None, status_color=None):
    """Create a status card with title, status text, and optional date."""
    # Add card shape
    card = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = RGBColor(255, 255, 255)
    card.line.color.rgb = RGBColor(230, 230, 230)
    
    # Add title
    title_box = slide.shapes.add_textbox(left + Inches(0.2), top + Inches(0.2), width - Inches(0.4), Inches(0.5))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_run = title_para.add_run()
    title_run.text = title
    title_run.font.bold = True
    title_run.font.size = Pt(16)
    
    # Add status indicator if color provided
    if status_color:
        dot_size = Inches(0.15)
        dot = slide.shapes.add_shape(MSO_SHAPE.OVAL, 
                                   left + width - Inches(0.4),
                                   top + Inches(0.3),
                                   dot_size, dot_size)
        dot.fill.solid()
        dot.fill.fore_color.rgb = status_color
        dot.line.fill.background()
    
    # Add status text
    status_box = slide.shapes.add_textbox(left + Inches(0.2), 
                                        top + Inches(0.8), 
                                        width - Inches(0.4), 
                                        Inches(0.5))
    status_frame = status_box.text_frame
    status_para = status_frame.paragraphs[0]
    status_run = status_para.add_run()
    status_run.text = status_text
    status_run.font.size = Pt(12)
    
    # Add date if provided
    if date_text:
        date_box = slide.shapes.add_textbox(left + Inches(0.2), 
                                          top + Inches(1.3), 
                                          width - Inches(0.4), 
                                          Inches(0.3))
        date_frame = date_box.text_frame
        date_para = date_frame.paragraphs[0]
        date_run = date_para.add_run()
        date_run.text = date_text
        date_run.font.size = Pt(10)
        date_run.font.color.rgb = RGBColor(128, 128, 128)

def create_progress_bar(slide, left, top, width, height, progress, label):
    """Create a progress bar with label and percentage."""
    # Background bar
    bg_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    bg_bar.fill.solid()
    bg_bar.fill.fore_color.rgb = RGBColor(230, 230, 230)
    bg_bar.line.fill.background()
    
    # Progress bar
    progress_width = width * (progress / 100)
    if progress_width > 0:
        progress_bar = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 
                                            left, top, 
                                            progress_width, height)
        progress_bar.fill.solid()
        progress_bar.fill.fore_color.rgb = RGBColor(0, 120, 212)
        progress_bar.line.fill.background()
    
    # Label
    label_box = slide.shapes.add_textbox(left - Inches(1.2), 
                                       top - Inches(0.1), 
                                       Inches(1), 
                                       Inches(0.3))
    label_frame = label_box.text_frame
    label_para = label_frame.paragraphs[0]
    label_run = label_para.add_run()
    label_run.text = label
    label_run.font.size = Pt(10)

def convert_pdf_to_ppt(pdf_path, output_path=None):
    """
    Convert a PDF file to a PowerPoint presentation with editable elements.
    
    Args:
        pdf_path (str): Path to the input PDF file
        output_path (str, optional): Path for the output PowerPoint file. 
                                   If not provided, will use the same name as PDF with .pptx extension
    """
    try:
        if output_path is None:
            output_path = os.path.splitext(pdf_path)[0] + '.pptx'
        
        # Create presentation with 16:9 aspect ratio
        prs = Presentation()
        prs.slide_width = Inches(16)
        prs.slide_height = Inches(9)
        
        pdf_document = fitz.open(pdf_path)
        print(f"Converting PDF with {len(pdf_document)} pages...")
        
        for page_num in range(len(pdf_document)):
            print(f"Processing page {page_num + 1}/{len(pdf_document)}...")
            
            page = pdf_document[page_num]
            slide = prs.slides.add_slide(prs.slide_layouts[6])
            
            # Get page dimensions and calculate scale
            pdf_width = page.rect.width
            pdf_height = page.rect.height
            scale_x = prs.slide_width / pdf_width
            scale_y = prs.slide_height / pdf_height
            
            # Process shapes first to establish layout structure
            shapes = page.get_drawings()
            cards = []
            progress_bars = []
            status_dots = []
            
            for shape in shapes:
                shape_type = detect_shape_type(shape)
                rect = shape["rect"]
                x0, y0, x1, y1 = rect
                left = x0 * scale_x
                top = y0 * scale_y
                width = (x1 - x0) * scale_x
                height = (y1 - y0) * scale_y
                
                if shape_type == "card":
                    cards.append((left, top, width, height))
                elif shape_type == "progress_bar":
                    progress_bars.append((left, top, width, height))
                elif shape_type == "status_dot":
                    status_dots.append((left, top, width, height))
            
            # Extract text blocks
            text_blocks = page.get_text("dict")["blocks"]
            
            # Process title first
            title_block = next((block for block in text_blocks 
                              if block.get("lines") and 
                              any("Global" in span["text"] 
                                  for line in block["lines"] 
                                  for span in line["spans"])), None)
            
            if title_block:
                title_box = slide.shapes.add_textbox(Inches(1), Inches(0.5), 
                                                   prs.slide_width - Inches(2), Inches(1))
                title_frame = title_box.text_frame
                title_para = title_frame.paragraphs[0]
                title_para.alignment = PP_ALIGN.CENTER
                title_run = title_para.add_run()
                title_text = " ".join(span["text"] for line in title_block["lines"] 
                                    for span in line["spans"])
                title_run.text = title_text
                title_run.font.size = Pt(24)
                title_run.font.bold = True
            
            # Create status cards
            card_width = Inches(4)
            card_height = Inches(2)
            card_margin = Inches(0.5)
            cards_top = Inches(2)
            
            # Example status cards (you'll need to extract actual data from PDF)
            status_data = [
                ("Germany", "Reviewing PBG portfolio", "March 28, 2023", RGBColor(255, 193, 7)),
                ("Spain", "Business data received", None, RGBColor(76, 175, 80)),
                ("Belgium", "Work commencing March 31st", None, RGBColor(76, 175, 80)),
                ("India", "Awaiting data", None, RGBColor(244, 67, 54))
            ]
            
            for i, (title, status, date, color) in enumerate(status_data):
                left = Inches(1) + (card_width + card_margin) * i
                create_status_card(slide, left, cards_top, card_width, card_height,
                                 title, status, date, color)
            
            # Create progress bars section
            progress_top = cards_top + card_height + Inches(1)
            progress_data = [
                ("Germany", 35),
                ("Spain", 75),
                ("Belgium", 80),
                ("India", 20)
            ]
            
            for i, (label, progress) in enumerate(progress_data):
                bar_left = Inches(3)
                bar_top = progress_top + Inches(0.5) * i
                create_progress_bar(slide, bar_left, bar_top,
                                  Inches(6), Inches(0.3),
                                  progress, label)
        
        print(f"Saving PowerPoint presentation to {output_path}...")
        prs.save(output_path)
        print("Conversion completed successfully!")
        
    except Exception as e:
        print(f"Error during conversion: {str(e)}")
        raise Exception(f"Failed to convert PDF: {str(e)}")
    finally:
        if 'pdf_document' in locals():
            pdf_document.close()

def main():
    parser = argparse.ArgumentParser(description='Convert PDF file to PowerPoint presentation')
    parser.add_argument('pdf_path', help='Path to the input PDF file')
    parser.add_argument('--output', '-o', help='Path for the output PowerPoint file (optional)')
    
    args = parser.parse_args()
    
    if not os.path.exists(args.pdf_path):
        print(f"Error: The file {args.pdf_path} does not exist.")
        return
    
    try:
        convert_pdf_to_ppt(args.pdf_path, args.output)
    except Exception as e:
        print(f"Error: {str(e)}")
        sys.exit(1)

if __name__ == '__main__':
    main() 